# Copyright (c) Microsoft. All rights reserved.
"""Read files the user shares in Teams and turn them into plain text.

When a user attaches a file, Teams sends a ``message`` activity with an
attachment of type ``application/vnd.microsoft.teams.file.download.info`` whose
``content.downloadUrl`` is a pre-authenticated URL. We download it, pull out the
text (plain text / code, PDF, DOCX, PPTX) and hand it straight to the model.
"""

from __future__ import annotations

import io
import json
import logging
from typing import Any

logger = logging.getLogger("github-copilot.files")

_FILE_DOWNLOAD_INFO = "application/vnd.microsoft.teams.file.download.info"
_MAX_CHARS = 20_000


async def read_shared_files(activity: Any) -> str:
    """Download every shared file and return its text, ready to inline in a prompt."""
    import httpx

    blocks: list[str] = []
    async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
        for att in getattr(activity, "attachments", None) or []:
            if (getattr(att, "content_type", "") or "").lower() != _FILE_DOWNLOAD_INFO:
                continue
            name = getattr(att, "name", None) or "attachment"
            url = _download_url(att)
            if not url:
                continue
            try:
                resp = await client.get(url)
                resp.raise_for_status()
                text = _extract_text(name, resp.content)[:_MAX_CHARS]
                blocks.append(f"--- {name} ---\n{text}")
            except Exception as ex:  # pylint: disable=broad-exception-caught
                logger.warning("Could not read %s: %s", name, ex)
    return "\n\n".join(blocks)


def _download_url(att: Any) -> str | None:
    """Return the pre-authenticated download URL for a Teams file attachment."""
    content = getattr(att, "content", None)
    if isinstance(content, str):
        try:
            content = json.loads(content)
        except (json.JSONDecodeError, ValueError):
            return None
    elif content is not None and not isinstance(content, dict):
        dump = getattr(content, "model_dump", None)
        content = dump(by_alias=True) if callable(dump) else None
    return content.get("downloadUrl") if isinstance(content, dict) else None


def _extract_text(name: str, data: bytes) -> str:
    """Extract plain text from ``data`` based on the file extension."""
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    if ext == "pdf":
        from pypdf import PdfReader

        return "\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(data)).pages)
    if ext == "docx":
        import docx

        return "\n".join(p.text for p in docx.Document(io.BytesIO(data)).paragraphs)
    if ext == "pptx":
        from pptx import Presentation

        return "\n".join(
            shape.text_frame.text
            for slide in Presentation(io.BytesIO(data)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
    return data.decode("utf-8", errors="replace")
