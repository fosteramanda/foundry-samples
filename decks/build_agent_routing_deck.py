"""Build "The Future of Human-Agent Interaction - Agents as Tools, Agents as Teammates".

Generates a 5-slide PowerPoint deck (16:9) that walks from today's model router,
through the agent toolbox, to the agent router, the autopilot router, and finally
delegation into an org-wide catalog of thousands of agents.

Usage:
    pip install python-pptx
    python decks/build_agent_routing_deck.py [output.pptx]
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- CoreAI-style palette -------------------------------------------------
INK = RGBColor(0x1B, 0x1B, 0x1F)        # primary text
MUTED = RGBColor(0x60, 0x5E, 0x5C)      # secondary text
HAIRLINE = RGBColor(0xE1, 0xDF, 0xDD)   # rules and card borders
CANVAS = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xFA, 0xF9, 0xF8)    # card fill
BLUE = RGBColor(0x0F, 0x6C, 0xBD)       # primary accent
BLUE_TINT = RGBColor(0xEB, 0xF3, 0xFC)
PURPLE = RGBColor(0x8A, 0x64, 0xD6)     # agents
PURPLE_TINT = RGBColor(0xF3, 0xEF, 0xFB)
TEAL = RGBColor(0x0E, 0x7C, 0x66)       # tools / grounding
TEAL_TINT = RGBColor(0xE8, 0xF4, 0xF1)
AMBER = RGBColor(0xB8, 0x6E, 0x00)      # identity / governance
AMBER_TINT = RGBColor(0xFD, 0xF4, 0xE3)

FONT = "Segoe UI"
DECK_TITLE = "The Future of Human-Agent Interaction"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)


# --- primitives -----------------------------------------------------------
def _set_text(frame, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    """runs: list of (text, size_pt, bold, color) tuples, one paragraph each."""
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(runs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        if spacing:
            para.space_after = Pt(spacing)
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    _set_text(box.text_frame, runs, align, anchor, spacing)
    return box


def card(slide, x, y, w, h, fill, line, radius=0.10, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = shadow
    shape.text_frame.word_wrap = True
    return shape


def chip(slide, x, y, w, h, label, sub=None, fill=SURFACE, line=HAIRLINE,
         label_color=INK, size=12, sub_size=9.5):
    shape = card(slide, x, y, w, h, fill, line)
    frame = shape.text_frame
    frame.margin_left = Inches(0.10)
    frame.margin_right = Inches(0.10)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)
    runs = [(label, size, True, label_color)]
    if sub:
        runs.append((sub, sub_size, False, MUTED))
    _set_text(frame, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shape


def arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    line_el = conn.line._get_or_add_ln()
    tail = line_el.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd", {}
    )
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    line_el.append(tail)
    if dashed:
        dash = line_el.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash", {}
        )
        dash.set("val", "dash")
        line_el.insert(0, dash)
    return conn


def center_bottom(shape):
    return shape.left + shape.width // 2, shape.top + shape.height


def center_top(shape):
    return shape.left + shape.width // 2, shape.top


def right_middle(shape):
    return shape.left + shape.width, shape.top + shape.height // 2


def left_middle(shape):
    return shape.left, shape.top + shape.height // 2


# --- slide chrome ---------------------------------------------------------
def new_slide(prs, eyebrow, title, subtitle, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CANVAS

    textbox(slide, MARGIN, Inches(0.48), Inches(11.8), Inches(0.28),
            [(eyebrow.upper(), 11, True, BLUE)])
    textbox(slide, MARGIN, Inches(0.78), Inches(11.8), Inches(0.55),
            [(title, 26, True, INK)])
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.42),
                                  Inches(1.1), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.fill.background()
    rule.shadow.inherit = False
    textbox(slide, MARGIN, Inches(1.60), Inches(11.0), Inches(0.35),
            [(subtitle, 14, False, MUTED)])

    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(6.82),
                                    Inches(11.83), Pt(1))
    footer.fill.solid()
    footer.fill.fore_color.rgb = HAIRLINE
    footer.line.fill.background()
    footer.shadow.inherit = False
    textbox(slide, MARGIN, Inches(6.95), Inches(9.0), Inches(0.3),
            [(DECK_TITLE, 9.5, False, MUTED)])
    textbox(slide, Inches(11.0), Inches(6.95), Inches(1.58), Inches(0.3),
            [(str(number), 9.5, False, MUTED)], align=PP_ALIGN.RIGHT)
    return slide


def takeaway(slide, text):
    bar = card(slide, MARGIN, Inches(6.05), Inches(11.83), Inches(0.62),
               BLUE_TINT, None, radius=0.08)
    bar.text_frame.margin_left = Inches(0.24)
    _set_text(bar.text_frame, [(text, 13, True, RGBColor(0x0A, 0x4A, 0x82))],
              anchor=MSO_ANCHOR.MIDDLE)
    return bar


def stage_label(slide, x, y, w, text, color=MUTED):
    return textbox(slide, x, y, w, Inches(0.24), [(text.upper(), 9.5, True, color)])


# --- slides ---------------------------------------------------------------
def slide1(prs):
    slide = new_slide(
        prs, "Where we are now",
        "How do I know which model to use?",
        "Foundry model router picks the right model for each request, so you stop hard-coding one model for every job.",
        1)

    top = Inches(2.35)
    prompt = chip(slide, MARGIN, top + Inches(0.85), Inches(2.3), Inches(0.95),
                  "Your request", "one prompt, no model choice", BLUE_TINT, BLUE, BLUE, 13)

    router = card(slide, Inches(3.75), top + Inches(0.35), Inches(2.9), Inches(1.95),
                  BLUE, None, radius=0.09)
    _set_text(router.text_frame,
              [("Foundry\nmodel router", 17, True, CANVAS),
               ("scores cost, latency,\nreasoning depth", 11, False,
                RGBColor(0xD9, 0xEA, 0xFA))],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=4)

    models = [
        ("Frontier reasoning", "hard, multi-step work"),
        ("Balanced general", "everyday requests"),
        ("Small + fast", "high-volume, low latency"),
        ("Domain / fine-tuned", "specialized answers"),
    ]
    mx, my, mh, gap = Inches(7.85), top - Inches(0.10), Inches(0.72), Inches(0.22)
    model_shapes = []
    for i, (name, sub) in enumerate(models):
        model_shapes.append(
            chip(slide, mx, my + i * (mh + gap), Inches(4.0), mh, name, sub,
                 SURFACE, HAIRLINE, INK, 12.5))

    stage_label(slide, MARGIN, top + Inches(0.55), Inches(2.3), "Input")
    stage_label(slide, Inches(3.75), top + Inches(0.05), Inches(2.9), "Routing", BLUE)
    stage_label(slide, mx, my - Inches(0.32), Inches(4.0), "Model catalog")

    arrow(slide, *right_middle(prompt), *left_middle(router), BLUE, 2)
    hub_x, hub_y = right_middle(router)
    for shape in model_shapes:
        arrow(slide, hub_x, hub_y, *left_middle(shape), MUTED, 1.25)

    takeaway(slide, "Solved once: the platform picks the model. The developer states intent, not infrastructure.")
    return slide


def slide2(prs):
    slide = new_slide(
        prs, "Where we are now",
        "How does my agent pick the best tool for the job?",
        "A Foundry agent plans over a governed toolbox and calls the right tool at each step.",
        2)

    agent = card(slide, MARGIN, Inches(2.55), Inches(3.2), Inches(2.35),
                 PURPLE_TINT, PURPLE, radius=0.08)
    _set_text(agent.text_frame,
              [("Foundry agent", 18, True, INK),
               ("instructions + memory + plan", 11.5, False, MUTED),
               ("1  understand the goal\n2  choose a tool\n3  observe the result\n4  repeat until done",
                11.5, False, INK)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=6)

    box = card(slide, Inches(5.35), Inches(2.30), Inches(7.23), Inches(3.05),
               SURFACE, HAIRLINE, radius=0.04)
    stage_label(slide, Inches(5.60), Inches(2.48), Inches(4.0), "Toolbox", TEAL)

    tools = [
        ("MCP servers", "connect any system"),
        ("Enterprise knowledge", "grounded retrieval"),
        ("Functions & APIs", "act in line-of-business apps"),
        ("Code interpreter", "compute and analysis"),
        ("Browser / computer use", "the long tail of UI"),
        ("Evaluations & tracing", "prove it worked"),
    ]
    tw, th, gx, gy = Inches(3.35), Inches(0.78), Inches(0.20), Inches(0.20)
    x0, y0 = Inches(5.60), Inches(2.82)
    for i, (name, sub) in enumerate(tools):
        col, row = i % 2, i // 2
        chip(slide, x0 + col * (tw + gx), y0 + row * (th + gy), tw, th,
             name, sub, CANVAS, HAIRLINE, INK, 12.5)

    arrow(slide, *right_middle(agent), left_middle(box)[0], right_middle(agent)[1],
          PURPLE, 2)
    textbox(slide, Inches(3.95), Inches(3.35), Inches(1.5), Inches(0.5),
            [("calls", 10.5, False, MUTED), ("returns", 10.5, False, MUTED)],
            align=PP_ALIGN.CENTER)

    takeaway(slide, "The agent already routes across tools. The next question is who routes across agents.")
    return slide


def slide3(prs):
    slide = new_slide(
        prs, "Where we are going next",
        "How do I pick the best agent for a task?",
        "An agent router treats sub-agents like tools - and every hop still runs on-behalf-of you.",
        3)

    user = chip(slide, MARGIN, Inches(3.05), Inches(2.05), Inches(1.05),
                "You", "one ask", BLUE_TINT, BLUE, BLUE, 14)

    router = card(slide, Inches(3.35), Inches(2.75), Inches(2.75), Inches(1.65),
                  PURPLE, None, radius=0.09)
    _set_text(router.text_frame,
              [("Agent router", 17, True, CANVAS),
               ("decomposes the goal,\npicks the right agent", 11, False,
                RGBColor(0xEA, 0xE2, 0xFB))],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=4)

    subs = [
        ("Research agent", "reads the corpus"),
        ("Finance agent", "runs the model"),
        ("Ticketing agent", "files the work"),
        ("Comms agent", "drafts the update"),
    ]
    sx, sy, sh, sgap = Inches(7.55), Inches(2.05), Inches(0.80), Inches(0.24)
    for i, (name, sub) in enumerate(subs):
        shape = chip(slide, sx, sy + i * (sh + sgap), Inches(3.35), sh,
                     name, sub, PURPLE_TINT, PURPLE, INK, 12.5)
        arrow(slide, *right_middle(router), *left_middle(shape), PURPLE, 1.25)
        badge = chip(slide, sx + Inches(3.45), sy + i * (sh + sgap) + Inches(0.16),
                     Inches(1.42), Inches(0.48), "on-behalf-of you", None,
                     AMBER_TINT, AMBER, AMBER, 9)
        badge.text_frame.word_wrap = True

    arrow(slide, *right_middle(user), *left_middle(router), BLUE, 2)

    note = card(slide, MARGIN, Inches(4.87), Inches(6.0), Inches(1.10),
                AMBER_TINT, AMBER, radius=0.06)
    note.text_frame.margin_left = Inches(0.18)
    note.text_frame.margin_right = Inches(0.18)
    _set_text(note.text_frame,
              [("What an agent identity actually changes", 12.5, True, INK),
               ("Every sub-agent inherits your permissions - it can never reach data you cannot. "
                "One consent, one audit trail, one place to revoke.", 11, False, MUTED)],
              anchor=MSO_ANCHOR.MIDDLE, spacing=3)

    takeaway(slide, "Agents become routable capabilities - without becoming a second, unmanaged set of credentials.")
    return slide


def slide4(prs):
    slide = new_slide(
        prs, "Where we are going after that",
        "A router that is proactive across my surfaces",
        "The autopilot works where I already am, notices the work, and delegates when it needs help.",
        4)

    surfaces = ["Teams", "Outlook", "Word & Excel", "IDE / repo", "Browser"]
    sw, sgap = Inches(2.13), Inches(0.20)
    sx0, sy = MARGIN, Inches(2.20)
    surface_shapes = []
    stage_label(slide, MARGIN, Inches(1.98), Inches(5.0), "Surfaces where work happens")
    for i, name in enumerate(surfaces):
        surface_shapes.append(
            chip(slide, sx0 + i * (sw + sgap), sy, sw, Inches(0.66), name, None,
                 SURFACE, HAIRLINE, INK, 12.5))

    pilot = card(slide, Inches(3.55), Inches(3.35), Inches(6.2), Inches(1.35),
                 BLUE, None, radius=0.09)
    _set_text(pilot.text_frame,
              [("Autopilot router", 18, True, CANVAS),
               ("persistent context  \u2022  proactive triggers  \u2022  works while you are away",
                11.5, False, RGBColor(0xD9, 0xEA, 0xFA))],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=4)

    span = pilot.width - Inches(1.2)
    for i, shape in enumerate(surface_shapes):
        target_x = pilot.left + Inches(0.6) + Emu(int(span * i / (len(surface_shapes) - 1)))
        arrow(slide, *center_bottom(shape), target_x, pilot.top, MUTED, 1.25, dashed=True)

    stage_label(slide, MARGIN, Inches(4.98), Inches(6.0), "Delegates a job-to-be-done", BLUE)
    jobs = [
        ("Close the quarter", "finance agent"),
        ("Ship the release notes", "comms agent"),
        ("Clear the backlog triage", "eng agent"),
    ]
    jw, jgap = Inches(3.6), Inches(0.24)
    for i, (job, owner) in enumerate(jobs):
        shape = chip(slide, MARGIN + i * (jw + jgap), Inches(5.22), jw, Inches(0.72),
                     job, owner, PURPLE_TINT, PURPLE, INK, 12.5)
        arrow(slide, center_bottom(pilot)[0], pilot.top + pilot.height,
              *center_top(shape), PURPLE, 1.25)

    takeaway(slide, "You stop starting every task. The autopilot starts it, and hands off what it should not do alone.")
    return slide


def slide5(prs):
    slide = new_slide(
        prs, "What about everything we already built",
        "We already have 1,000s of agents. Now what?",
        "Nothing is thrown away - the autopilot becomes the front door to agents you already published.",
        5)

    pilot = card(slide, MARGIN, Inches(2.05), Inches(3.15), Inches(1.50),
                 BLUE, None, radius=0.09)
    _set_text(pilot.text_frame,
              [("Autopilot router", 16.5, True, CANVAS),
               ("one relationship\nwith the employee", 11, False,
                RGBColor(0xD9, 0xEA, 0xFA))],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=4)

    catalog = card(slide, Inches(5.20), Inches(2.05), Inches(7.38), Inches(3.05),
                   SURFACE, HAIRLINE, radius=0.04)
    stage_label(slide, Inches(5.39), Inches(2.22), Inches(5.0),
                "Agent Store  \u2022  My agents  \u2022  1,000s already deployed", TEAL)

    tiles = [
        ("HR onboarding", "Contoso HR"),
        ("Expense policy", "Finance"),
        ("Security triage", "SecOps"),
        ("Field sales brief", "Sales ops"),
        ("Supplier lookup", "Procurement"),
        ("Legal redline", "Legal"),
        ("+ 1,000s more", "every team, already shipped"),
        ("Partner agents", "ISV catalog"),
    ]
    tw, th, gx, gy = Inches(3.4), Inches(0.62), Inches(0.18), Inches(0.16)
    x0, y0 = Inches(5.39), Inches(2.58)
    for i, (name, owner) in enumerate(tiles):
        col, row = i % 2, i // 2
        chip(slide, x0 + col * (tw + gx), y0 + row * (th + gy), tw, th,
             name, owner, CANVAS, HAIRLINE, INK, 12)

    arrow(slide, *right_middle(pilot), *left_middle(catalog), BLUE, 2)
    textbox(slide, Inches(3.95), Inches(2.62), Inches(1.3), Inches(0.3),
            [("delegates", 10.5, False, MUTED)], align=PP_ALIGN.CENTER)

    gov = card(slide, MARGIN, Inches(3.80), Inches(4.05), Inches(1.30),
               AMBER_TINT, AMBER, radius=0.06)
    gov.text_frame.margin_left = Inches(0.18)
    gov.text_frame.margin_right = Inches(0.18)
    _set_text(gov.text_frame,
              [("Same identity fabric", 12.5, True, INK),
               ("Each agent keeps its own identity, owner and policy. "
                "Delegation is logged, scoped and revocable.", 11, False, MUTED)],
              anchor=MSO_ANCHOR.MIDDLE, spacing=3)

    takeaway(slide, "Autopilots do not replace your agents - they make the thousands you already have reachable from one conversation.")
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for maker in (slide1, slide2, slide3, slide4, slide5):
        maker(prs)
    prs.save(path)
    return path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "Agents-as-Tools-Agents-as-Teammates.pptx"
    print(build(out))
